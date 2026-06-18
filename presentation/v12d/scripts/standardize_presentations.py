#!/usr/bin/env python3
"""Apply focused branding and layout fixes to the V12d thesis decks."""

from __future__ import annotations

import argparse
import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


EMU_PER_INCH = 914400
CALIBRI = "Calibri"
NAVY = RGBColor(0x12, 0x26, 0x3A)
SLATE = RGBColor(0x55, 0x69, 0x78)
SUBTITLE_SLATE = RGBColor(0x64, 0x74, 0x8B)
TEAL = RGBColor(0x02, 0x80, 0x90)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_geometry(shape, x: float, y: float, width: float, height: float) -> None:
    shape.left = Inches(x)
    shape.top = Inches(y)
    shape.width = Inches(width)
    shape.height = Inches(height)


def find_text_shape(slide, exact_text: str):
    for shape in slide.shapes:
        if getattr(shape, "text", "").strip() == exact_text:
            return shape
    raise ValueError(f"Text shape not found: {exact_text!r}")


def set_text(
    shape,
    text: str,
    *,
    size: float,
    color: RGBColor,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.MIDDLE,
    word_wrap: bool = True,
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.10)
    frame.margin_right = Inches(0.10)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.05)
    frame.word_wrap = word_wrap
    frame.vertical_anchor = anchor

    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = CALIBRI
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def set_picture_alt_text(picture, description: str) -> None:
    c_nv_pr = picture._element.nvPicPr.cNvPr
    c_nv_pr.set("descr", description)
    c_nv_pr.set("title", "VGU and h_da institutional logos")


def standardize_title_slide(prs: Presentation, logo_path: Path) -> None:
    slide = prs.slides[0]

    header = find_text_shape(slide, "MASTER THESIS · ITS 2024 · 21024003")
    affiliation = find_text_shape(
        slide, "Hochschule Darmstadt · Fachbereich Informatik"
    )
    title = find_text_shape(slide, "Multimodal RAG for Leishmaniasis Diagnosis")
    subtitle = find_text_shape(
        slide,
        "Traceable evidence, model-dependent retrieval, and safety-gated demonstration",
    )
    author = find_text_shape(
        slide,
        "Nguyen Tuan Ngoc · Supervisor: Prof. Dr. Andreas Müller · "
        "Co-Supervisor: Dr. Tran Duc Khanh · May 2026",
    )

    set_geometry(header, 0.40, 0.22, 6.75, 0.30)
    set_text(header, header.text, size=9, color=SLATE)

    set_geometry(affiliation, 0.40, 0.57, 6.75, 0.38)
    set_text(
        affiliation,
        "Computer Science (M.Sc.) · HDA / VGU",
        size=11,
        color=SLATE,
    )

    set_geometry(title, 0.40, 1.30, 6.80, 1.95)
    set_text(
        title,
        title.text,
        size=32,
        color=NAVY,
        bold=True,
        anchor=MSO_ANCHOR.TOP,
    )

    set_geometry(subtitle, 0.40, 3.55, 8.90, 0.55)
    set_text(
        subtitle,
        subtitle.text,
        size=14,
        color=SLATE,
        italic=True,
    )

    set_geometry(author, 0.40, 4.45, 8.65, 0.90)
    set_text(author, author.text, size=13, color=SLATE, bold=True)

    for shape in slide.shapes:
        if (
            not getattr(shape, "text", "")
            and abs(shape.left / EMU_PER_INCH - 0.40) < 0.02
            and abs(shape.top / EMU_PER_INCH - 4.35) < 0.02
        ):
            set_geometry(shape, 0.40, 4.30, 5.50, 0.03)
            break

    logo = slide.shapes.add_picture(
        str(logo_path),
        Inches(7.32),
        Inches(0.08),
        width=Inches(2.24),
        height=Inches(1.51),
    )
    set_picture_alt_text(
        logo,
        "Vietnamese-German University and Hochschule Darmstadt "
        "(Darmstadt University of Applied Sciences) logos.",
    )


def standardize_architecture_slide(prs: Presentation) -> None:
    slide = prs.slides[3]
    title = find_text_shape(slide, "Two-lane retrieval architecture")
    subtitle = find_text_shape(
        slide,
        "Text retrieval and image-linked retrieval are fused, gated, and audited "
        "before presentation.",
    )
    page = find_text_shape(slide, "04")
    caption = find_text_shape(
        slide,
        "Architecture figure reused from the thesis pipeline diagram; the demo "
        "exposes retrieval, generation, uncertainty gate, and final route.",
    )
    picture = next(shape for shape in slide.shapes if shape.shape_type == 13)

    set_geometry(title, 0.40, 0.15, 9.00, 0.55)
    set_text(title, title.text, size=24, color=NAVY, bold=True)

    set_geometry(subtitle, 0.40, 0.72, 9.00, 0.28)
    set_text(
        subtitle,
        subtitle.text,
        size=11,
        color=SUBTITLE_SLATE,
        italic=True,
    )

    set_geometry(page, 9.40, 0.10, 0.40, 0.30)
    page.fill.solid()
    page.fill.fore_color.rgb = TEAL
    page.line.fill.background()
    set_text(
        page,
        "04",
        size=8,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )

    # Preserve the architecture image's aspect ratio while creating clean
    # separation between the subtitle and bottom caption.
    target_height = 4.32
    aspect_ratio = picture.width / picture.height
    target_width = target_height * aspect_ratio
    set_geometry(
        picture,
        (10.0 - target_width) / 2.0,
        1.02,
        target_width,
        target_height,
    )

    set_geometry(caption, 0.55, 5.42, 8.90, 0.14)
    set_text(
        caption,
        caption.text,
        size=5.5,
        color=SUBTITLE_SLATE,
        align=PP_ALIGN.CENTER,
        word_wrap=False,
    )


def explicit_calibri_for_evaluation_section(prs: Presentation) -> None:
    """Make the generated evaluation/appendix section use the core deck font."""
    if len(prs.slides) <= 14:
        return
    for slide_index in range(14, len(prs.slides)):
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = CALIBRI


def is_top_right_page_shape(shape) -> bool:
    return (
        shape.left / EMU_PER_INCH >= 9.30
        and shape.top / EMU_PER_INCH <= 0.50
        and shape.width / EMU_PER_INCH <= 0.60
        and shape.height / EMU_PER_INCH <= 0.50
    )


def is_bottom_right_page_shape(shape) -> bool:
    return (
        shape.left / EMU_PER_INCH >= 8.80
        and shape.top / EMU_PER_INCH >= 5.00
        and getattr(shape, "text", "").strip().isdigit()
    )


def is_closing_slide(slide) -> bool:
    texts = {
        getattr(shape, "text", "").strip()
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    }
    return "Thank You" in texts and "Questions?" in texts


def remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def clone_page_badge(slide, templates, slide_number: int) -> None:
    next_shape_id = max((shape.shape_id for shape in slide.shapes), default=1) + 1
    for template, contains_number in templates:
        element = copy.deepcopy(template)
        c_nv_pr = element.find(".//" + qn("p:cNvPr"))
        if c_nv_pr is not None:
            c_nv_pr.set("id", str(next_shape_id))
            c_nv_pr.set("name", f"Page badge {slide_number:02d}")
        next_shape_id += 1

        if contains_number:
            text_nodes = element.findall(".//" + qn("a:t"))
            if text_nodes:
                text_nodes[0].text = f"{slide_number:02d}"
                for text_node in text_nodes[1:]:
                    text_node.text = ""

        slide.shapes._spTree.insert_element_before(element, "p:extLst")


def unify_page_badges(prs: Presentation) -> None:
    """Apply one top-right page-number format to every content slide."""
    template_slide = prs.slides[1]
    template_shapes = [
        shape for shape in template_slide.shapes if is_top_right_page_shape(shape)
    ]
    if len(template_shapes) != 2:
        raise ValueError(
            "Expected the slide 2 page badge to contain one background shape "
            "and one number shape"
        )
    templates = [
        (copy.deepcopy(shape._element), bool(getattr(shape, "text", "").strip()))
        for shape in template_shapes
    ]

    for slide_index in range(1, len(prs.slides)):
        slide = prs.slides[slide_index]
        for shape in list(slide.shapes):
            if is_top_right_page_shape(shape) or is_bottom_right_page_shape(shape):
                remove_shape(shape)

        if slide_index >= 14:
            for shape in slide.shapes:
                if (
                    abs(shape.left / EMU_PER_INCH - 0.55) < 0.02
                    and abs(shape.top / EMU_PER_INCH - 0.32) < 0.02
                    and abs(shape.height / EMU_PER_INCH - 0.24) < 0.02
                ):
                    shape.width = Inches(8.65)
                    break

        if is_closing_slide(slide):
            continue

        clone_page_badge(slide, templates, slide_index + 1)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("input_pptx", type=Path)
    parser.add_argument("output_pptx", type=Path)
    parser.add_argument("--logo", required=True, type=Path)
    args = parser.parse_args()

    prs = Presentation(str(args.input_pptx))
    if len(prs.slides) not in {14, 21, 22, 25, 31}:
        raise ValueError(
            "Expected a 14-, 21-, 22-, 25-, or 31-slide V12d deck, "
            f"found {len(prs.slides)} slides"
        )
    if (prs.slide_width / EMU_PER_INCH, prs.slide_height / EMU_PER_INCH) != (
        10.0,
        5.625,
    ):
        raise ValueError("Unexpected slide size; expected 10 x 5.625 inches")

    standardize_title_slide(prs, args.logo)
    standardize_architecture_slide(prs)
    explicit_calibri_for_evaluation_section(prs)
    unify_page_badges(prs)

    args.output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.output_pptx))
    print(f"Saved {args.output_pptx}")


if __name__ == "__main__":
    main()
