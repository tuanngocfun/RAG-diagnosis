#!/usr/bin/env python3
"""Validate V12d deck structure, provenance, and exported PDF text."""

from __future__ import annotations

import json
import os
import subprocess
import sys
import unicodedata
import zipfile
from pathlib import Path

from pptx import Presentation


PROJECT_ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = PROJECT_ROOT / "FINAL_PRESENTATION.pptx"
PDF_PATH = PROJECT_ROOT / "FINAL_PRESENTATION.pdf"
RESULT_DIR = PROJECT_ROOT / "data" / "heldout_evaluation_results"
COMPARISON_PATH = (
    PROJECT_ROOT / "data" / "gemma4_rag_norag_comparison" / "comparison_summary.json"
)
TRACE_PATH = (
    PROJECT_ROOT / "data" / "exact_rag_trace_appendix" / "trace_summary.json"
)
RECAPTURE_MANIFEST = PROJECT_ROOT / "re-capture" / "latest_live_recapture_manifest.json"
EMU_PER_INCH = 914400
DEFAULT_SPLIT_ROOT = Path(
    "/home/ngocnt/Leishmaniasis_v3/rag/instructions/process/14/"
    "details_analysis/rtx_titan/structured_cases_v4/"
    "leishmaniasis_verified_v2"
)
SPLIT_ROOT = Path(os.environ.get("LEISH_SPLIT_ROOT", DEFAULT_SPLIT_ROOT))
TRAIN_JSONL = (
    SPLIT_ROOT
    / "nonleish_additions"
    / "generated"
    / "train_phase1b_tierAB.jsonl"
)
TEST_JSONL = SPLIT_ROOT / "test_p14_v7_normalized.jsonl"
EVAL_QUERY_JSONL = SPLIT_ROOT / "eval_queries_p14_v7_mixed56.jsonl"

CASE_CHECKS = {
    16: {
        "case_id": "PMC7516301_01",
        "live_rank_type": "MCL",
        "required": ("19-year-old", "Syria", "throat"),
        "forbidden": (),
        "appendix": 20,
    },
    17: {
        "case_id": "PMC7456484_01",
        "live_rank_type": "MCL",
        "required": ("37-year-old", "Bhutan", "nasal"),
        "forbidden": ("34-year-old male",),
        "appendix": 21,
    },
    18: {
        "case_id": "PMC10026180_04",
        "live_rank_type": "CL",
        "required": ("11-year-old", "ataxia telangiectasia", "nose"),
        "forbidden": (
            "18-" + "month-old",
            "peri" + "anal",
            "abs" + "cess",
            "chronic ulcer" + "ation",
        ),
        "appendix": 22,
    },
}

FORBIDDEN_OVERCLAIMS = (
    "rag " + "doubles",
    "smoking " + "gun",
    "rag " + "essential",
    "+34 " + "percentage points",
    "three-case accuracy",
)


def normalize(text: str) -> str:
    normalized = unicodedata.normalize("NFKC", text).casefold()
    normalized = normalized.translate(str.maketrans("‐‑‒–—−", "------"))
    return " ".join(normalized.split())


def slide_text(slide) -> str:
    return "\n".join(
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    )


def appendix_contains_clinical_text(appendix_text: str, clinical_text: str) -> bool:
    appendix = normalize(appendix_text)
    paragraphs = [p for p in clinical_text.split("\n\n") if p.strip()]
    if len(paragraphs) <= 1:
        return normalize(clinical_text) in appendix
    return all(normalize(paragraph) in appendix for paragraph in paragraphs)


def jsonl_case_counts(path: Path) -> dict[str, int]:
    counts: dict[str, int] = {}
    with path.open("r", encoding="utf-8") as handle:
        for line in handle:
            if not line.strip():
                continue
            case_id = json.loads(line).get("case_id", "")
            counts[case_id] = counts.get(case_id, 0) + 1
    return counts


def load_live_recapture_manifest() -> dict:
    if not RECAPTURE_MANIFEST.exists():
        raise FileNotFoundError(f"Missing live recapture manifest: {RECAPTURE_MANIFEST}")
    manifest = json.loads(RECAPTURE_MANIFEST.read_text(encoding="utf-8"))
    output_dir = Path(manifest["output_dir"])
    if not output_dir.exists():
        raise FileNotFoundError(f"Missing live recapture directory: {output_dir}")
    return manifest


def live_recapture_dir() -> Path:
    return Path(load_live_recapture_manifest()["output_dir"])


def load_live_case(case_id: str) -> dict:
    path = live_recapture_dir() / f"{case_id}_live_gpu_result.json"
    if not path.exists():
        raise FileNotFoundError(f"Missing live recapture case result: {path}")
    return json.loads(path.read_text(encoding="utf-8"))


def live_fields(case_id: str) -> dict:
    live = load_live_case(case_id)
    raw = live.get("raw_response") or {}
    runtime = raw.get("runtime_metadata") or {}
    rank = dict(live.get("rank_fields") or {})
    return {
        "request_id": raw.get("request_id", ""),
        "provider_mode": raw.get("provider_mode", ""),
        "model_name": raw.get("model_name", ""),
        "elapsed_seconds": live.get("elapsed_seconds", ""),
        "query_image_tensor_count": runtime.get("query_image_tensor_count", ""),
        "safety_state": raw.get("safety_state", ""),
        "rank1": rank.get("rank1", ""),
        "rank1_type": rank.get("rank1_type", ""),
        "rank1_confidence": rank.get("rank1_confidence", ""),
        "evidence": raw.get("evidence") or [],
    }


def validate_provenance(errors: list[str]) -> None:
    for path in (TRAIN_JSONL, TEST_JSONL, EVAL_QUERY_JSONL):
        if not path.exists():
            errors.append(f"Missing split artifact: {path}")
    if errors:
        return

    train_counts = jsonl_case_counts(TRAIN_JSONL)
    test_counts = jsonl_case_counts(TEST_JSONL)
    eval_counts = jsonl_case_counts(EVAL_QUERY_JSONL)
    if sum(train_counts.values()) != 121:
        errors.append(
            f"Official retrieval corpus must have 121 rows, found "
            f"{sum(train_counts.values())}"
        )

    for check in CASE_CHECKS.values():
        case_id = check["case_id"]
        if train_counts.get(case_id, 0) != 0:
            errors.append(f"{case_id} appears in the official retrieval corpus")
        if test_counts.get(case_id, 0) != 1:
            errors.append(f"{case_id} held-out count is not one")
        if eval_counts.get(case_id, 0) != 1:
            errors.append(f"{case_id} eval-query count is not one")


def validate_pptx(errors: list[str]) -> Presentation:
    with zipfile.ZipFile(PPTX_PATH) as package:
        bad_member = package.testzip()
        if bad_member:
            errors.append(f"PPTX ZIP integrity failed at {bad_member}")

    prs = Presentation(str(PPTX_PATH))
    if len(prs.slides) != 31:
        errors.append(f"Expected 31 slides, found {len(prs.slides)}")

    width = prs.slide_width / EMU_PER_INCH
    height = prs.slide_height / EMU_PER_INCH
    if (width, height) != (10.0, 5.625):
        errors.append(f"Unexpected slide size: {width} x {height}")

    for slide_number, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            left = shape.left / EMU_PER_INCH
            top = shape.top / EMU_PER_INCH
            right = (shape.left + shape.width) / EMU_PER_INCH
            bottom = (shape.top + shape.height) / EMU_PER_INCH
            if left < -0.01 or top < -0.01 or right > 10.01 or bottom > 5.635:
                errors.append(
                    f"Slide {slide_number}: shape {shape.shape_id} is outside "
                    f"the canvas ({left:.2f}, {top:.2f}, {right:.2f}, {bottom:.2f})"
                )

    overview = normalize(slide_text(prs.slides[14]))
    for required in (
        "selected example outcomes",
        "aggregate thesis metrics remain the benchmark",
        "illustrative, not an accuracy sample",
        "official 121-case experimental retrieval corpus",
        "local defense demo kb",
        "silver labels are pipeline references",
        "fresh gemma 4 live-demo outputs",
        "label conflict",
        "evidence-attribution",
    ):
        if normalize(required) not in overview:
            errors.append(f"Slide 15 is missing required wording: {required}")
    if "performance summary" in overview:
        errors.append("Slide 15 still contains 'Performance " + "Summary'")

    all_text = normalize("\n".join(slide_text(slide) for slide in prs.slides))
    for stale in (
        "18-" + "month-old",
        "peri" + "anal abs" + "cess",
        "34-year-old " + "male",
        "reviewer assessment; not gemma output",
        "expected diagnosis",
        "safety: generated_support",
        "three-case accuracy",
        "non-leish specificity",
        "specificity example",
        "rank 1 remained non-leishmaniasis",
        "non-leish -> non-leish",
        *FORBIDDEN_OVERCLAIMS,
    ):
        if normalize(stale) in all_text:
            errors.append(f"Deck contains stale wording: {stale}")

    slide17_text = normalize(slide_text(prs.slides[16]))
    if "model rank 1 mcl" not in slide17_text:
        errors.append("Slide 17 no longer preserves the PKDL -> MCL subtype story")
    if "family-level" in slide17_text or "false non-leish" in slide17_text:
        errors.append("Slide 17 contains v11b-style family-level failure wording")

    slide18_text = normalize(slide_text(prs.slides[17]))
    for required in (
        "label-conflict stress test",
        "verified non-leish / pseudolabel cl",
        "model rank 1 cl",
        "leishmaniasis-plausible",
        "not specificity proof",
    ):
        if required not in slide18_text:
            errors.append(f"Slide 18 is missing label-conflict wording: {required}")

    for result_slide, check in CASE_CHECKS.items():
        result = json.loads(
            (RESULT_DIR / f"{check['case_id']}_result.json").read_text(
                encoding="utf-8"
            )
        )
        result_text = normalize(slide_text(prs.slides[result_slide - 1]))
        appendix_text = normalize(slide_text(prs.slides[check["appendix"] - 1]))
        clinical_text = normalize(result["clinical_text"])
        provenance = result.get("split_provenance", {})

        if result_slide == 18:
            if "label conflict" not in result_text:
                errors.append("Slide 18 does not label the reference conflict")
        elif "silver reference label" not in result_text:
            errors.append(
                f"Slide {result_slide} does not label the reference as silver"
            )
        if "local defense demo kb" not in result_text:
            errors.append(
                f"Slide {result_slide} does not disclose the runtime demo KB"
            )
        live = live_fields(check["case_id"])
        if live["rank1_type"] != check["live_rank_type"]:
            errors.append(
                f"Live recapture rank type for {check['case_id']} is "
                f"{live['rank1_type']}, expected {check['live_rank_type']}"
            )
        for required_live in (
            live["rank1_type"],
            live["rank1_confidence"],
            live["rank1"],
        ):
            if normalize(str(required_live)) not in result_text:
                errors.append(
                    f"Slide {result_slide} is missing live recapture field "
                    f"for {check['case_id']}: {required_live}"
                )
        if not provenance.get(
            "clinical_retrieval_corpus_source", ""
        ).endswith("train_phase1b_tierAB.jsonl"):
            errors.append(
                f"{check['case_id']} does not reference the canonical "
                "121-case train artifact"
            )
        if not provenance.get("runtime_retrieval_kb_source"):
            errors.append(
                f"{check['case_id']} does not record the runtime demo KB"
            )
        if result.get("label_contract", {}).get(
            "ground_truth_status"
        ) != "silver_reference_only":
            errors.append(
                f"{check['case_id']} does not preserve the silver-label contract"
            )

        for required in check["required"]:
            term = normalize(required)
            if term not in result_text:
                errors.append(
                    f"Slide {result_slide} is missing source anchor: {required}"
                )
            if term not in clinical_text:
                errors.append(
                    f"{check['case_id']} source JSON is missing anchor: {required}"
                )
        for forbidden in check["forbidden"]:
            if normalize(forbidden) in result_text:
                errors.append(
                    f"Slide {result_slide} contains stale term: {forbidden}"
                )
        if not appendix_contains_clinical_text(
            slide_text(prs.slides[check["appendix"] - 1]),
            result["clinical_text"],
        ):
            errors.append(
                f"Slide {check['appendix']} does not preserve the full clinical "
                f"text for {check['case_id']}"
            )

    validate_comparison_slides(prs, errors)
    validate_trace_appendix_slides(prs, errors)

    return prs


def validate_comparison_slides(prs: Presentation, errors: list[str]) -> None:
    if not COMPARISON_PATH.exists():
        errors.append(f"Missing comparison artifact: {COMPARISON_PATH}")
        return
    comparison = json.loads(COMPARISON_PATH.read_text(encoding="utf-8"))
    if comparison.get("comparison_label") != "official Gemma 4 experiment-pipeline comparison":
        errors.append("Comparison artifact does not use the official Gemma 4 source label")
    if ("med" + "gemma") in normalize(json.dumps(comparison, ensure_ascii=False)):
        errors.append("Comparison artifact unexpectedly references Med" + "Gemma")
    if "not used because the backend still returned retrieved support chunks" not in normalize(
        comparison.get("supersedes", "")
    ):
        errors.append("Comparison artifact does not supersede the failed backend no-RAG attempt")

    expected_ids = ("PMC7516301_01", "PMC7456484_01", "PMC10026180_04")
    cases = comparison.get("cases", [])
    if tuple(case.get("case_id") for case in cases) != expected_ids:
        errors.append("Comparison artifact case order is not the expected three held-out IDs")

    for index, case_id in enumerate(expected_ids, start=23):
        text = normalize(slide_text(prs.slides[index - 1]))
        for required in (
            case_id,
            "official gemma 4 experiment-pipeline comparison",
            "no-rag condition",
            "rag condition",
            "q&a backup only",
            "aggregate thesis metrics remain the benchmark",
        ):
            if normalize(required) not in text:
                errors.append(f"Slide {index} is missing comparison wording: {required}")

    slide23 = normalize(slide_text(prs.slides[22]))
    if "rank 1 type: non-leishmaniasis" not in slide23 or "rank 1 type: mcl" not in slide23:
        errors.append("Slide 23 does not show the intended no-RAG/RAG case-1 contrast")
    slide24 = normalize(slide_text(prs.slides[23]))
    if "subtype-resolution challenge" not in slide24:
        errors.append("Slide 24 does not preserve the PKDL subtype-challenge framing")
    slide25 = normalize(slide_text(prs.slides[24]))
    for required in (
        "label-conflict",
        "evidence-attribution stress test",
        "not specificity proof",
    ):
        if required not in slide25:
            errors.append(f"Slide 25 is missing case-3 audit framing: {required}")
    if "both conditions remain non-leishmaniasis, supporting the specificity" in slide25:
        errors.append("Slide 25 still presents case 3 as specificity proof")


def validate_trace_appendix_slides(prs: Presentation, errors: list[str]) -> None:
    if not TRACE_PATH.exists():
        errors.append(f"Missing exact RAG trace artifact: {TRACE_PATH}")
        return

    trace = json.loads(TRACE_PATH.read_text(encoding="utf-8"))
    if trace.get("version") != "V12d":
        errors.append("Exact trace artifact does not identify V12d")
    run = trace.get("official_rag_run", {})
    if run.get("retriever_method") != "hybrid":
        errors.append("Official RAG trace does not record retriever_method=hybrid")
    if run.get("rerank") is not True or run.get("summary_rerank") is not True:
        errors.append("Official RAG trace does not record rerank=true")
    if run.get("retrieval_top_k") != 20:
        errors.append("Official RAG trace does not record retrieval_top_k=20")

    expected_ids = ("PMC7516301_01", "PMC7456484_01", "PMC10026180_04")
    cases = trace.get("cases", [])
    if tuple(case.get("case_id") for case in cases) != expected_ids:
        errors.append("Exact trace case order is not the expected three held-out IDs")
        return

    for pair_index, case in enumerate(cases):
        official_slide_num = 26 + pair_index * 2
        fresh_slide_num = official_slide_num + 1
        official_text = normalize(slide_text(prs.slides[official_slide_num - 1]))
        fresh_text = normalize(slide_text(prs.slides[fresh_slide_num - 1]))

        case_id = case["case_id"]
        qid = case["qid"]
        for required in (
            case_id,
            qid,
            "official gemma 4 experiment-pipeline trace",
            "rerank-enabled final context list used for generation",
            "hybrid",
            "true",
            "20",
        ):
            if normalize(required) not in official_text:
                errors.append(
                    f"Slide {official_slide_num} is missing trace wording: {required}"
                )

        for context in case["official_rag_trace"]["top_contexts_for_slide"][:3]:
            for field in (
                context["doc_id"],
                str(context["score"]),
                context["diagnosis_type"],
                context["label_source"],
            ):
                if normalize(field) not in official_text:
                    errors.append(
                        f"Slide {official_slide_num} is missing exact context "
                        f"field for {case_id}: {field}"
                    )
            context_prefix = " ".join(context["text_prefix_260"].split())[:90]
            if normalize(context_prefix) not in official_text:
                errors.append(
                    f"Slide {official_slide_num} is missing exact context "
                    f"excerpt prefix for {case_id}: {context['doc_id']}"
                )

        fields = live_fields(case_id)
        for field in (
            case_id,
            fields["request_id"],
            fields["model_name"],
            fields["provider_mode"],
            str(fields["elapsed_seconds"]),
            str(fields["query_image_tensor_count"]),
            fields["safety_state"],
            fields["rank1"],
            fields["rank1_type"],
            fields["rank1_confidence"],
        ):
            if normalize(field) not in fresh_text:
                errors.append(
                    f"Slide {fresh_slide_num} is missing exact fresh GPU field "
                    f"for {case_id}: {field}"
                )

        for evidence in fields["evidence"][:3]:
            for field in (
                evidence.get("chunk_id", ""),
                str(evidence.get("score", "")),
                evidence.get("title", ""),
                evidence.get("diagnosis_label", ""),
                str(evidence.get("confirmatory", "")),
            ):
                if normalize(field) not in fresh_text:
                    errors.append(
                        f"Slide {fresh_slide_num} is missing exact evidence "
                        f"field for {case_id}: {field}"
                    )
            excerpt_prefix = " ".join(str(evidence.get("excerpt", "")).split())[:90]
            if normalize(excerpt_prefix) not in fresh_text:
                errors.append(
                    f"Slide {fresh_slide_num} is missing exact evidence "
                    f"excerpt prefix for {case_id}: {evidence.get('chunk_id', '')}"
                )


def validate_pdf(errors: list[str]) -> None:
    pdf_text = ""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        try:
            result = subprocess.run(
                ["pdftotext", str(PDF_PATH), "-"],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
            )
            pdf_text = result.stdout
        except (subprocess.CalledProcessError, FileNotFoundError) as exc:
            errors.append(f"PDF text validation could not run: {exc}")
            return
    else:
        reader = PdfReader(str(PDF_PATH))
        if len(reader.pages) != 31:
            errors.append(f"Expected 31 PDF pages, found {len(reader.pages)}")
        pdf_text = "\n".join(page.extract_text() or "" for page in reader.pages)

    text = normalize(pdf_text)
    for required in (
        "live demo case inspection",
        "selected example outcomes",
        "silver reference label",
        "label conflict",
        "verified non-leish / pseudolabel cl",
        "leishmaniasis-plausible",
        "local defense demo kb",
        "11-year-old",
        "ataxia telangiectasia",
        "official gemma 4 experiment-pipeline comparison",
        "official gemma 4 experiment-pipeline trace",
        "fresh real-gpu output audit",
        "rerank-enabled final context list used for generation",
        "q&a backup only",
    ):
        if normalize(required) not in text:
            errors.append(f"PDF is missing required wording: {required}")
    for stale in (
        "18-" + "month-old",
        "peri" + "anal abs" + "cess",
        "performance summary",
        "expected diagnosis",
        "safety: generated_support",
        "specificity example",
        "rank 1 remained non-leishmaniasis",
        *FORBIDDEN_OVERCLAIMS,
    ):
        if normalize(stale) in text:
            errors.append(f"PDF contains stale wording: {stale}")


def main() -> int:
    errors: list[str] = []
    for path in (PPTX_PATH, PDF_PATH):
        if not path.exists():
            errors.append(f"Missing artifact: {path}")

    if not errors:
        validate_provenance(errors)
    if not errors:
        validate_pptx(errors)
        validate_pdf(errors)

    if errors:
        print("V12d validation failed:", file=sys.stderr)
        for error in errors:
            print(f"- {error}", file=sys.stderr)
        return 1

    print(
        "V12d validation passed: 31 slides/pages, silver-label and "
        "corpus/runtime provenance aligned, comparison and exact-trace "
        "appendices bounded."
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
