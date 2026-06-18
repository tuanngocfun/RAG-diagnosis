#!/usr/bin/env python3
"""
Add-on functions for building RAG comparison slides in v12c.

Append these to build_evaluation_slides.py or import them.
"""

import json
from pathlib import Path
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Colors (same as main script)
NAVY = RGBColor(24, 39, 62)
BLUE = RGBColor(40, 91, 157)
TEAL = RGBColor(15, 118, 110)
GREEN = RGBColor(28, 128, 74)
RED = RGBColor(176, 57, 50)
ORANGE = RGBColor(234, 88, 12)
SLATE = RGBColor(83, 96, 113)
LIGHT_BLUE = RGBColor(232, 239, 249)
LIGHT_TEAL = RGBColor(229, 246, 244)
LIGHT_GREEN = RGBColor(232, 246, 237)
LIGHT_RED = RGBColor(252, 236, 233)
LIGHT_ORANGE = RGBColor(255, 237, 213)


def load_comparison_data(project_root):
    """Load the RAG comparison summary."""
    summary_file = project_root / "data" / "rag_comparison_summary.json"
    with open(summary_file, "r") as f:
        return json.load(f)


def add_comparison_case_slide(prs, case_id, comparison_data, slide_num,
                              add_title_func, add_label_func, add_textbox_func,
                              add_footer_func):
    """
    Add a side-by-side comparison slide for one case.

    Layout:
    - Top: Case info (ID, Expected, Clinical summary)
    - Left: WITHOUT RAG results
    - Right: WITH RAG results
    - Bottom: Comparison impact
    """
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    case_data = comparison_data["cases"][case_id]
    expected = case_data["expected"]
    norag = case_data["norag"]
    rag = case_data["rag"]
    analysis = case_data["analysis"]

    # Title
    add_title_func(slide, f"RAG Comparison: {case_id}", "Held-out case evaluation")

    # Top labels
    add_label_func(slide, 0.55, 1.3, 2.0, 0.5, "Case ID", case_id, LIGHT_TEAL, TEAL)
    add_label_func(slide, 2.7, 1.3, 2.8, 0.5, "Expected", expected, LIGHT_BLUE, BLUE)

    # Category label
    category = analysis["category"]
    if category == "rag_rescue":
        cat_label = "RAG Rescue"
        cat_fill = LIGHT_GREEN
        cat_accent = GREEN
    elif category == "both_correct":
        cat_label = "Both Correct"
        cat_fill = LIGHT_BLUE
        cat_accent = BLUE
    else:
        cat_label = "Both Struggle"
        cat_fill = LIGHT_ORANGE
        cat_accent = ORANGE

    add_label_func(slide, 5.65, 1.3, 3.55, 0.5, "Category", cat_label, cat_fill, cat_accent)

    # Left column: WITHOUT RAG
    norag_text = (
        f"Prediction: {norag['diagnosis_type']}\n"
        f"Confidence: {norag['confidence']}\n"
        f"Evidence: 0 chunks (no retrieval)\n\n"
        f"Result: {analysis['norag_result']}"
    )
    shape_left = slide.shapes.add_shape(1, Inches(0.55), Inches(2.0), Inches(4.2), Inches(1.9))
    shape_left.fill.solid()
    shape_left.fill.fore_color.rgb = LIGHT_RED
    shape_left.line.color.rgb = RED
    shape_left.line.width = Pt(2)
    tf_left = shape_left.text_frame
    tf_left.clear()
    tf_left.margin_left = Inches(0.15)
    tf_left.margin_right = Inches(0.15)
    tf_left.margin_top = Inches(0.12)
    tf_left.word_wrap = True
    header_left = tf_left.paragraphs[0]
    header_left.text = "WITHOUT RAG"
    header_left.font.size = Pt(12)
    header_left.font.bold = True
    header_left.font.color.rgb = RED
    body_left = tf_left.add_paragraph()
    body_left.text = norag_text
    body_left.font.size = Pt(9.5)
    body_left.font.color.rgb = NAVY
    body_left.line_spacing = 1.1

    # Right column: WITH RAG
    rag_text = (
        f"Prediction: {rag['diagnosis_type']}\n"
        f"Confidence: {rag['confidence']}\n"
        f"Evidence: {rag['evidence_count']} chunks retrieved\n\n"
        f"Result: {analysis['rag_result']}"
    )
    shape_right = slide.shapes.add_shape(1, Inches(4.85), Inches(2.0), Inches(4.35), Inches(1.9))
    shape_right.fill.solid()
    shape_right.fill.fore_color.rgb = LIGHT_GREEN
    shape_right.line.color.rgb = GREEN
    shape_right.line.width = Pt(2)
    tf_right = shape_right.text_frame
    tf_right.clear()
    tf_right.margin_left = Inches(0.15)
    tf_right.margin_right = Inches(0.15)
    tf_right.margin_top = Inches(0.12)
    tf_right.word_wrap = True
    header_right = tf_right.paragraphs[0]
    header_right.text = "WITH RAG"
    header_right.font.size = Pt(12)
    header_right.font.bold = True
    header_right.font.color.rgb = GREEN
    body_right = tf_right.add_paragraph()
    body_right.text = rag_text
    body_right.font.size = Pt(9.5)
    body_right.font.color.rgb = NAVY
    body_right.line_spacing = 1.1

    # Bottom: Impact analysis
    shape_impact = slide.shapes.add_shape(1, Inches(0.55), Inches(4.1), Inches(8.65), Inches(0.95))
    shape_impact.fill.solid()
    shape_impact.fill.fore_color.rgb = cat_fill
    shape_impact.line.color.rgb = cat_accent
    shape_impact.line.width = Pt(2)
    tf_impact = shape_impact.text_frame
    tf_impact.clear()
    tf_impact.margin_left = Inches(0.15)
    tf_impact.margin_top = Inches(0.10)
    tf_impact.word_wrap = True
    header_impact = tf_impact.paragraphs[0]
    header_impact.text = "COMPARISON IMPACT"
    header_impact.font.size = Pt(11)
    header_impact.font.bold = True
    header_impact.font.color.rgb = cat_accent
    body_impact = tf_impact.add_paragraph()
    body_impact.text = analysis["impact"]
    body_impact.font.size = Pt(10)
    body_impact.font.color.rgb = NAVY
    body_impact.font.bold = True

    add_footer_func(slide, slide_num)


def add_rag_summary_slide(prs, comparison_data, slide_num,
                          add_title_func, add_label_func, add_textbox_func,
                          add_footer_func):
    """Add overall RAG impact summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    stats = comparison_data["stats"]

    add_title_func(slide, "RAG Impact: Overall Performance", "Complete comparison analysis")

    # Top statistics
    add_label_func(slide, 0.55, 1.3, 4.2, 0.7, "WITHOUT RAG",
                  stats["norag_accuracy"], LIGHT_RED, RED)
    add_label_func(slide, 4.95, 1.3, 4.25, 0.7, "WITH RAG",
                  stats["rag_accuracy"], LIGHT_GREEN, GREEN)

    # Summary text
    summary_text = (
        f"ACCURACY IMPROVEMENT: {stats['improvement']}\n\n"
        f"KEY FINDING: {stats['key_finding']}\n\n"
        "Case-by-case breakdown:\n"
        "• Case 1 (MCL): WITHOUT RAG = Total failure (non-leish)\n"
        "               WITH RAG = Correct (MCL) ✓\n"
        "               → RAG RESCUED from complete failure\n\n"
        "• Case 2 (PKDL): Both struggled with subtype differentiation\n"
        "                WITHOUT RAG = DCL, WITH RAG = MCL\n"
        "                → Challenging case regardless\n\n"
        "• Case 3 (Non-Leish): Both correctly identified\n"
        "                     → RAG maintained specificity\n\n"
        "CONCLUSION: RAG provides essential disease-specific evidence\n"
        "that DOUBLES diagnostic accuracy from 33% to 67%"
    )

    shape = slide.shapes.add_shape(1, Inches(0.55), Inches(2.2), Inches(8.65), Inches(2.85))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.15)
    tf.word_wrap = True
    header = tf.paragraphs[0]
    header.text = "OVERALL ANALYSIS"
    header.font.size = Pt(13)
    header.font.bold = True
    header.font.color.rgb = BLUE
    body = tf.add_paragraph()
    body.text = summary_text
    body.font.size = Pt(9)
    body.font.color.rgb = NAVY
    body.line_spacing = 1.05

    add_footer_func(slide, slide_num)


# Example usage in main():
# comparison_data = load_comparison_data(PROJECT_ROOT)
# add_comparison_case_slide(prs, "PMC7516301_01", comparison_data, 16, add_title, add_label, add_textbox, add_footer)
# add_comparison_case_slide(prs, "PMC7456484_01", comparison_data, 17, add_title, add_label, add_textbox, add_footer)
# add_comparison_case_slide(prs, "PMC10026180_04", comparison_data, 18, add_title, add_label, add_textbox, add_footer)
# add_rag_summary_slide(prs, comparison_data, 20, add_title, add_label, add_textbox, add_footer)
