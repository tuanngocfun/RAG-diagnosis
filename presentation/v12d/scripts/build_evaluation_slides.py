#!/usr/bin/env python3
"""Build the V12d held-out case inspection and Q&A appendix slides."""

from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
EVAL_DIR = PROJECT_ROOT / "data" / "heldout_evaluation_results"
COMPARISON_DIR = PROJECT_ROOT / "data" / "gemma4_rag_norag_comparison"
COMPARISON_SUMMARY = COMPARISON_DIR / "comparison_summary.json"
TRACE_DIR = PROJECT_ROOT / "data" / "exact_rag_trace_appendix"
TRACE_SUMMARY = TRACE_DIR / "trace_summary.json"
RECAPTURE_DIR = PROJECT_ROOT / "re-capture"
RECAPTURE_MANIFEST = RECAPTURE_DIR / "latest_live_recapture_manifest.json"
LOGO_PATH = PROJECT_ROOT / "assets" / "logos" / "vgu_hda_lockup.png"
SOURCE_PPTX = (
    PROJECT_ROOT / "source" / "thesis_defense_slides_all_light_v6_harnessed.pptx"
)
OUTPUT_PPTX = (
    PROJECT_ROOT / "source" / "thesis_defense_slides_v5_with_evaluation.pptx"
)

CASE_PRESENTATION = {
    "PMC7516301_01": (
        "19-year-old man from Syria with 6-month throat irritation progressing "
        "to severe throat pain and fever; exam noted prior arm scar and "
        "oropharyngeal/laryngeal mucosal lesions.",
        ("19-year-old", "Syria", "throat"),
        (),
    ),
    "PMC7456484_01": (
        "37-year-old woman from eastern Bhutan with extensive facial plaques, "
        "eyelid damage/scarring, nasal mucosal extension, and prior treated "
        "visceral leishmaniasis history.",
        ("37", "Bhutan", "nasal"),
        ("34-year-old male",),
    ),
    "PMC10026180_04": (
        "11-year-old girl with ataxia telangiectasia, recurrent respiratory "
        "and ear infections, low IgG, and a spreading nose rash extending to "
        "the nostrils and left cheek; the image/text are leishmaniasis-plausible "
        "but the pipeline labels conflict.",
        ("11-year-old", "ataxia telangiectasia", "nose"),
        (
            "18-" + "month-old",
            "peri" + "anal",
            "abs" + "cess",
            "chronic ulcer" + "ation",
        ),
    ),
}

CORE_WORDING_REPLACEMENTS = {
    "What v6 was missing: core result map":
        "Core result map: evidence behind the main claims",
    "v6 was compact for a live supervisor demo; thesis_v44b-1 contains a richer "
    "result chain that needs a backup map.":
        "The spoken deck is intentionally compact; this map links headline "
        "claims to the supporting thesis results.",
    "Why v6 felt short": "Why the spoken deck is compressed",
    "It compressed many result families into slides 5-10 to preserve a "
    "30-minute defense flow.":
        "The spoken deck compresses many result families into slides 5-10 to "
        "preserve a 30-minute defense flow.",
    "What v3 adds": "What the backup map adds",
    "Use this slide when supervisors ask where each result family lives in "
    "thesis_v44b-1.":
        "Use this map when supervisors ask how each result family connects to "
        "the thesis results chapter and appendix.",
    "These are the missing mechanism checks behind the short v6 RAG story.":
        "These mechanism checks explain why the observed RAG effect is "
        "conditional.",
    "This slide collects appendix evidence that supports the demo story without "
    "overclaiming deployment readiness.":
        "This slide collects appendix evidence that supports the defense "
        "narrative without overclaiming deployment readiness.",
    "These are the missing core result families behind the short v6 narrative.":
        "These result families support the defense narrative without implying "
        "deployment readiness.",
}

# Colors
NAVY = RGBColor(24, 39, 62)
BLUE = RGBColor(40, 91, 157)
TEAL = RGBColor(15, 118, 110)
GREEN = RGBColor(28, 128, 74)
RED = RGBColor(176, 57, 50)
ORANGE = RGBColor(234, 88, 12)
SLATE = RGBColor(83, 96, 113)
LIGHT_BLUE = RGBColor(232, 239, 249)
LIGHT_TEAL = RGBColor(229, 246, 244)
LIGHT_GREEN = RGBColor(232, 246, 237)
LIGHT_RED = RGBColor(252, 236, 233)
LIGHT_ORANGE = RGBColor(255, 237, 213)
WHITE = RGBColor(255, 255, 255)


def set_picture_alt_text(picture, description: str) -> None:
    c_nv_pr = picture._element.nvPicPr.cNvPr
    c_nv_pr.set("descr", description)
    c_nv_pr.set("title", "VGU and h_da institutional logos")


def clear_text_frame(shape) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE


def add_textbox(slide, x, y, w, h, text, size=14, color=NAVY, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    clear_text_frame(box)
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.line_spacing = 1.08
    if align:
        p.alignment = align
    return box


def add_title(slide, title: str, kicker: str) -> None:
    add_textbox(slide, 0.55, 0.32, 8.9, 0.24, kicker.upper(), 8.5, TEAL, True)
    add_textbox(slide, 0.55, 0.58, 8.9, 0.42, title, 23, NAVY, True)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.08), Inches(8.9), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()


def add_label(slide, x, y, w, h, label, value, fill, accent) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label.upper()
    p.font.size = Pt(7.5)
    p.font.color.rgb = accent
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(11)
    p2.font.color.rgb = NAVY
    p2.font.bold = True


def add_bullet_box(slide, x, y, w, h, title, bullets, fill, accent, bullet_size=9.5) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True
    header = tf.paragraphs[0]
    header.text = title
    header.font.size = Pt(12)
    header.font.bold = True
    header.font.color.rgb = accent
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(bullet_size)
        p.font.color.rgb = NAVY
        p.level = 0
        p.line_spacing = 1.05


def add_text_panel(slide, x, y, w, h, title, body, fill, accent, body_size=8.6) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    tf.word_wrap = True
    header = tf.paragraphs[0]
    header.text = title
    header.font.size = Pt(12)
    header.font.bold = True
    header.font.color.rgb = accent
    body_lines = body.split("\n")
    for line in body_lines:
        body_paragraph = tf.add_paragraph()
        body_paragraph.text = line
        body_paragraph.font.size = Pt(body_size)
        body_paragraph.font.color.rgb = NAVY
        body_paragraph.line_spacing = 1.0
        body_paragraph.space_after = Pt(0)


def add_assessment_box(slide, x, y, w, h, bullets, fill, accent) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.09)
    tf.margin_bottom = Inches(0.08)
    tf.word_wrap = True

    header = tf.paragraphs[0]
    header.text = "LLM Council Review Synthesis"
    header.font.size = Pt(11.5)
    header.font.bold = True
    header.font.color.rgb = accent
    header.alignment = PP_ALIGN.CENTER

    source = tf.add_paragraph()
    source.text = "Review synthesis; not Gemma output"
    source.font.size = Pt(8.6)
    source.font.bold = True
    source.font.color.rgb = SLATE
    source.alignment = PP_ALIGN.CENTER
    source.line_spacing = 0.9

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(8.25)
        p.font.color.rgb = NAVY
        p.line_spacing = 0.98


def add_picture_fit(slide, image_path: Path, x, y, max_w, max_h) -> None:
    with Image.open(image_path) as image:
        width_px, height_px = image.size
    image_ratio = width_px / height_px
    box_ratio = max_w / max_h
    if image_ratio >= box_ratio:
        width = max_w
        height = max_w / image_ratio
    else:
        height = max_h
        width = max_h * image_ratio
    x_offset = x + (max_w - width) / 2
    y_offset = y + (max_h - height) / 2
    slide.shapes.add_picture(
        str(image_path),
        Inches(x_offset),
        Inches(y_offset),
        width=Inches(width),
        height=Inches(height),
    )


def add_footer(slide, slide_number: int) -> None:
    add_textbox(slide, 0.55, 5.28, 5.5, 0.18,
                "Live case inspection - fresh local-GPU demo output",
                7.5, SLATE)
    page = add_textbox(slide, 8.95, 5.28, 0.5, 0.18, str(slide_number), 7.5, SLATE)
    page.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_comparison_footer(slide, slide_number: int) -> None:
    add_textbox(
        slide,
        0.55,
        5.28,
        6.7,
        0.18,
        "Q&A appendix only - official Gemma 4 experiment-pipeline comparison",
        7.5,
        SLATE,
    )
    page = add_textbox(slide, 8.95, 5.28, 0.5, 0.18, str(slide_number), 7.5, SLATE)
    page.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_trace_footer(slide, slide_number: int) -> None:
    add_textbox(
        slide,
        0.55,
        5.28,
        7.8,
        0.18,
        "Q&A appendix only - exact RAG trace and fresh real-GPU audit evidence",
        7.5,
        SLATE,
    )
    page = add_textbox(slide, 8.95, 5.28, 0.5, 0.18, str(slide_number), 7.5, SLATE)
    page.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_closing_slide(prs: Presentation) -> None:
    """Add the unnumbered closing slide shown during questions."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    accent = slide.shapes.add_shape(
        1, Inches(0.0), Inches(0.0), Inches(0.18), Inches(5.625)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    add_textbox(
        slide,
        0.55,
        0.55,
        5.7,
        0.3,
        "MASTER THESIS DEFENSE",
        10,
        SLATE,
        True,
    )
    add_textbox(slide, 0.55, 1.22, 5.8, 0.82, "Thank You", 36, NAVY, True)
    add_textbox(slide, 0.55, 2.12, 5.8, 0.58, "Questions?", 24, TEAL, True)
    add_textbox(
        slide,
        0.55,
        2.72,
        7.2,
        0.24,
        "Feasible and traceable under constraints; not deployment-ready.",
        12,
        SLATE,
        True,
    )

    line = slide.shapes.add_shape(
        1, Inches(0.55), Inches(3.12), Inches(5.5), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()

    add_textbox(
        slide,
        0.55,
        3.42,
        8.7,
        0.48,
        "Multimodal RAG for Leishmaniasis Diagnosis",
        15,
        NAVY,
        True,
    )
    add_textbox(
        slide,
        0.55,
        4.10,
        8.7,
        0.34,
        "Nguyen Tuan Ngoc · Computer Science (M.Sc.)",
        12,
        SLATE,
        True,
    )
    add_textbox(
        slide,
        0.55,
        4.54,
        8.7,
        0.3,
        "Vietnamese-German University · Hochschule Darmstadt (h_da)",
        10,
        SLATE,
    )

    logo = slide.shapes.add_picture(
        str(LOGO_PATH),
        Inches(7.32),
        Inches(0.08),
        width=Inches(2.24),
        height=Inches(1.51),
    )
    set_picture_alt_text(
        logo,
        "Vietnamese-German University and Hochschule Darmstadt "
        "(Darmstadt University of Applied Sciences) logos.",
    )


def add_case_overview_slide(prs: Presentation, slide_num: int) -> None:
    """Overview slide introducing the three evaluation cases."""
    live_fields = {
        case_id: live_case_fields(case_id)
        for case_id in ("PMC7516301_01", "PMC7456484_01", "PMC10026180_04")
    }
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(
        slide,
        "Live Demo Case Inspection",
        "Fresh local-GPU multimodal evaluation",
    )

    # Summary table
    add_label(slide, 0.55, 1.3, 2.8, 0.62, "Case 1: PMC7516301_01",
              f"Silver: MCL | Live Rank 1: {live_fields['PMC7516301_01']['rank1_type']}",
              LIGHT_GREEN, GREEN)
    add_label(slide, 3.5, 1.3, 2.8, 0.62, "Case 2: PMC7456484_01",
              f"Silver: PKDL | Live Rank 1: {live_fields['PMC7456484_01']['rank1_type']}",
              LIGHT_ORANGE, ORANGE)
    add_label(slide, 6.25, 1.3, 2.95, 0.62, "Case 3: PMC10026180_04",
              f"Label conflict: verified Non-Leish / pseudolabel CL | Live Rank 1: {live_fields['PMC10026180_04']['rank1_type']}",
              LIGHT_RED, RED)

    # Overview
    add_bullet_box(
        slide, 0.55, 2.08, 8.65, 1.22,
        "Case Inspection Setup",
        [
            "Three selected queries from the 56-case held-out set; illustrative, not an accuracy sample",
            "Verified absent from the official 121-case experimental retrieval corpus",
            "Fresh Gemma 4 live-demo outputs used the local defense demo KB on NVIDIA TITAN RTX (4-bit)",
            "Silver labels are pipeline references; aggregate thesis metrics remain the benchmark"
        ],
        LIGHT_BLUE, BLUE, bullet_size=9.0
    )

    # Results summary
    add_bullet_box(
        slide, 0.55, 3.45, 4.2, 1.45,
        "Selected Example Outcomes",
        [
            "Case 1 (MCL): live rank-1 type matches the silver reference",
            "Case 2 (PKDL): live rank-1 stays in the leishmaniasis family but shifts to MCL",
            "Case 3: verified Non-Leish / pseudolabel CL; live output stays leishmaniasis-family"
        ],
        LIGHT_GREEN, GREEN, bullet_size=9.2
    )

    add_bullet_box(
        slide, 4.9, 3.45, 4.3, 1.45,
        "Key Findings",
        [
            "Inspectable input-output-evidence traces under local GPU constraints",
            "Case 2 exposes subtype-resolution limits",
            "Case 3 exposes label-conflict and evidence-attribution limits",
            "These cases do not establish clinical validity or deployment readiness",
            "Images provide context, not diagnosis; full inputs: slides 20-22"
        ],
        LIGHT_TEAL, TEAL, bullet_size=8.6
    )

    add_footer(slide, slide_num)


def load_case_result(case_id: str) -> dict:
    """Load case result JSON."""
    result_file = EVAL_DIR / f"{case_id}_result.json"
    with open(result_file, "r", encoding="utf-8") as f:
        return json.load(f)


def load_live_recapture_manifest() -> dict:
    if not RECAPTURE_MANIFEST.exists():
        raise FileNotFoundError(
            "Live demo recapture manifest missing. "
            "Run scripts/recapture_live_demo_cases.py first."
        )
    with RECAPTURE_MANIFEST.open("r", encoding="utf-8") as handle:
        manifest = json.load(handle)
    output_dir = Path(manifest["output_dir"])
    if not output_dir.exists():
        raise FileNotFoundError(f"Live demo recapture directory missing: {output_dir}")
    return manifest


def latest_live_recapture_dir() -> Path:
    return Path(load_live_recapture_manifest()["output_dir"])


def load_live_case_result(case_id: str) -> dict:
    result_file = latest_live_recapture_dir() / f"{case_id}_live_gpu_result.json"
    if not result_file.exists():
        raise FileNotFoundError(f"Live demo recapture result missing: {result_file}")
    with result_file.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def live_case_fields(case_id: str) -> dict:
    live = load_live_case_result(case_id)
    fields = dict(live["rank_fields"])
    raw = live["raw_response"]
    runtime = raw.get("runtime_metadata") or {}
    fields.update(
        {
            "elapsed_seconds": live["elapsed_seconds"],
            "provider_mode": raw.get("provider_mode", "unknown"),
            "model_name": raw.get("model_name", "unknown"),
            "request_id": raw.get("request_id", ""),
            "safety_state": raw.get("safety_state", "unknown"),
            "query_image_tensor_count": runtime.get("query_image_tensor_count", "unknown"),
            "source_path": str(latest_live_recapture_dir() / f"{case_id}_live_gpu_result.json"),
        }
    )
    return fields


def format_seconds(value) -> str:
    try:
        return f"{float(value):.1f}s"
    except (TypeError, ValueError):
        return f"{value}s"


def reference_label_for_slide(case_id: str, result: dict) -> tuple[str, str]:
    if case_id == "PMC10026180_04":
        return "Label Conflict", "Verified Non-Leish / pseudolabel CL"
    return (
        "Silver Reference Label",
        result.get("silver_reference_label", result["expected_diagnosis"]),
    )


def evidence_for_slide(evidence: list[dict], limit: int = 3) -> list[dict]:
    rows = []
    for rank, item in enumerate(evidence[:limit], start=1):
        rows.append(
            {
                "rank": rank,
                "chunk_id": str(item.get("chunk_id", "")),
                "score": item.get("score", ""),
                "title": str(item.get("title", "")),
                "diagnosis_label": str(item.get("diagnosis_label", "")),
                "confirmatory": item.get("confirmatory", ""),
                "excerpt": str(item.get("excerpt", "")),
            }
        )
    return rows


def attach_live_recapture_audits(trace_summary: dict) -> dict:
    """Replace stale fresh-audit appendix fields with latest live recapture."""
    for case in trace_summary.get("cases", []):
        case_id = case["case_id"]
        live = load_live_case_result(case_id)
        raw = live["raw_response"]
        fields = live_case_fields(case_id)
        case["fresh_gpu_audit"] = {
            "source_dir": str(latest_live_recapture_dir()),
            "source_path": fields["source_path"],
            "result": live,
            "fields": {
                "request_id": fields["request_id"],
                "provider_mode": fields["provider_mode"],
                "model_name": fields["model_name"],
                "rank1": fields["rank1"],
                "rank1_type": fields["rank1_type"],
                "rank1_confidence": fields["rank1_confidence"],
                "safety_state": fields["safety_state"],
                "elapsed_seconds": fields["elapsed_seconds"],
                "query_image_tensor_count": fields["query_image_tensor_count"],
            },
            "evidence_for_slide": evidence_for_slide(raw.get("evidence") or []),
        }
    return trace_summary


def normalize_text(text: str) -> str:
    return " ".join(text.casefold().split())


def validated_case_summary(case_id: str, result: dict) -> str:
    summary, required_terms, forbidden_terms = CASE_PRESENTATION[case_id]
    clinical_text = normalize_text(result["clinical_text"])
    summary_text = normalize_text(summary)

    missing = [
        term for term in required_terms
        if normalize_text(term) not in clinical_text
        or normalize_text(term) not in summary_text
    ]
    stale = [
        term for term in forbidden_terms
        if normalize_text(term) in summary_text
    ]
    if missing or stale:
        raise ValueError(
            f"{case_id}: presentation summary validation failed; "
            f"missing={missing}, stale={stale}"
        )
    return summary


def refine_core_wording(prs: Presentation) -> None:
    for slide in list(prs.slides)[:14]:
        for shape in slide.shapes:
            current = getattr(shape, "text", "").strip()
            replacement = CORE_WORDING_REPLACEMENTS.get(current)
            if not replacement:
                continue
            paragraph = shape.text_frame.paragraphs[0]
            if len(paragraph.runs) != 1:
                raise ValueError(f"Cannot safely replace formatted text: {current!r}")
            paragraph.runs[0].text = replacement


def extract_rank1(markdown_text: str) -> dict:
    """Extract rank 1 diagnosis from model output."""
    lines = markdown_text.split("\n")
    result = {
        "diagnosis": "Unknown",
        "diagnosis_type": "Unknown",
        "confidence": "Unknown"
    }

    for line in lines:
        if "**Rank 1 supportive consideration:**" in line:
            result["diagnosis"] = line.split("**Rank 1 supportive consideration:**")[1].strip()
        elif "**Rank 1 Diagnosis Type:**" in line:
            result["diagnosis_type"] = line.split("**Rank 1 Diagnosis Type:**")[1].strip()
        elif "**Rank 1 Confidence:**" in line:
            result["confidence"] = line.split("**Rank 1 Confidence:**")[1].strip()

    return result


def humanize_output_route(safety_state: str) -> str:
    routes = {
        "generated_support": "retrieval-supported response",
        "generated_no_support": "generated without retrieved support",
        "abstained": "abstained",
    }
    return routes.get(safety_state, safety_state.replace("_", " "))


def add_case_result_slide(prs: Presentation, case_id: str, slide_num: int,
                          category: str) -> None:
    """Add slide showing one case's input and output with fixed layout."""
    result = load_case_result(case_id)
    live = load_live_case_result(case_id)
    live_response = live["raw_response"]
    live_fields = live_case_fields(case_id)
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Layout constants - two column design with proper spacing
    LEFT_COL_X = 0.55
    LEFT_COL_WIDTH = 4.3
    GUTTER = 0.2
    RIGHT_COL_X = LEFT_COL_X + LEFT_COL_WIDTH + GUTTER  # 5.05
    RIGHT_COL_WIDTH = 4.15

    ROW1_Y = 1.25
    ROW1_HEIGHT = 0.55
    ROW2_Y = ROW1_Y + ROW1_HEIGHT + 0.15  # 1.95
    ROW2_HEIGHT = 1.35
    ROW3_Y = ROW2_Y + ROW2_HEIGHT + 0.15  # 3.45
    ROW3_HEIGHT = 0.95
    ROW4_Y = ROW3_Y + ROW3_HEIGHT + 0.1  # 4.5
    ROW4_HEIGHT = 0.75

    # Title color based on category
    if category == "success":
        title_text = f"Case {slide_num - 15}: {case_id} — Concordant Example"
        color_fill = LIGHT_GREEN
        color_accent = GREEN
    elif category == "limitation":
        title_text = f"Case {slide_num - 15}: {case_id} — Subtype Mismatch"
        color_fill = LIGHT_ORANGE
        color_accent = ORANGE
    elif category == "label_conflict":
        title_text = f"Case {slide_num - 15}: {case_id} — Label-Conflict Stress Test"
        color_fill = LIGHT_RED
        color_accent = RED
    else:
        title_text = f"Case {slide_num - 15}: {case_id}"
        color_fill = LIGHT_BLUE
        color_accent = BLUE

    add_title(slide, title_text, "Fresh local-GPU live demo output")

    # Top row: silver reference vs model output (spans full width)
    model_output = {
        "diagnosis": live_fields["rank1"],
        "diagnosis_type": live_fields["rank1_type"],
        "confidence": live_fields["rank1_confidence"],
    }
    reference_heading, reference_label = reference_label_for_slide(case_id, result)
    add_label(slide, LEFT_COL_X, ROW1_Y, 2.5, ROW1_HEIGHT,
              reference_heading, reference_label,
              color_fill, color_accent)
    add_label(slide, LEFT_COL_X + 2.65, ROW1_Y, 2.7, ROW1_HEIGHT, "Model Rank 1",
              model_output["diagnosis_type"], color_fill, color_accent)
    add_label(slide, RIGHT_COL_X + 1.85, ROW1_Y, 1.15, ROW1_HEIGHT, "Confidence",
              model_output["confidence"], LIGHT_BLUE, BLUE)
    add_label(slide, RIGHT_COL_X + 3.1, ROW1_Y, 1.05, ROW1_HEIGHT, "Time",
              format_seconds(live_fields["elapsed_seconds"]), LIGHT_TEAL, TEAL)

    # LEFT COLUMN
    # Row 2 Left: Clinical input
    appendix_slide = slide_num + 4
    clinical_summary = validated_case_summary(case_id, result)
    add_bullet_box(
        slide, LEFT_COL_X, ROW2_Y, LEFT_COL_WIDTH, ROW2_HEIGHT,
        "Clinical Input (Blinded)",
        [
            clinical_summary,
            f"Held-out case absent from the official 121-case experimental retrieval corpus; full text is on slide {appendix_slide}.",
            "Fresh live output used the local defense demo KB; target and confirmatory fields were withheld.",
        ],
        LIGHT_BLUE, BLUE, bullet_size=7.9
    )

    # Row 3 Left: Model output summary
    output_bullets = [
        f"Top diagnosis: {model_output['diagnosis']}",
        f"Runtime support: {len(result.get('evidence', []))} retrieved chunks",
        f"Output route: {humanize_output_route(result.get('safety_state', 'unknown'))}"
    ]
    add_bullet_box(
        slide, LEFT_COL_X, ROW3_Y, LEFT_COL_WIDTH, ROW3_HEIGHT,
        "Model Output Summary",
        [
            f"Top diagnosis: {model_output['diagnosis']}",
            f"Runtime support: {len(live_response.get('evidence', []))} retrieved chunks",
            f"Output route: {humanize_output_route(live_response.get('safety_state', 'unknown'))}"
        ],
        LIGHT_TEAL, TEAL, bullet_size=9
    )

    # RIGHT COLUMN
    # Row 2 Right: Image (with proper boundary)
    image_path = Path(live["image_path"])
    if not image_path.exists():
        image_path = EVAL_DIR / f"{case_id}_input.png"
    if image_path.exists():
        # Calculate image dimensions to fit within boundary
        img_width = RIGHT_COL_WIDTH * 0.9  # Leave some margin
        img_height = ROW2_HEIGHT * 0.85  # Leave space for label

        add_picture_fit(slide, image_path, RIGHT_COL_X, ROW2_Y, img_width, img_height)
        add_textbox(slide, RIGHT_COL_X, ROW2_Y + img_height + 0.05,
                   RIGHT_COL_WIDTH, 0.25, "Input image", 8, SLATE, bold=True)

    # Row 3 Right: Performance assessment (separate from image)
    if case_id == "PMC7516301_01":
        assessment = [
            "Reviews agree MCL Rank 1 is a strong concordant positive example",
            "Image is visually consistent/contextual, not diagnostic proof"
        ]
    elif case_id == "PMC7456484_01":
        assessment = [
            "Reviews agree the family is retained but PKDL is not named",
            "Demo KB lacks a PKDL anchor; subtype limitation"
        ]
    elif case_id == "PMC10026180_04":
        assessment = [
            "Reviews flag a label conflict; full-text case is leishmaniasis-plausible",
            "Live CL/MCL-family output is not specificity proof; confirmation required"
        ]
    else:
        assessment = ["No assessment available"]

    add_assessment_box(
        slide,
        RIGHT_COL_X,
        ROW3_Y,
        RIGHT_COL_WIDTH,
        ROW3_HEIGHT,
        assessment,
        color_fill,
        color_accent,
    )

    # Bottom row: Evidence (full width)
    if live_response.get('evidence'):
        evidence_text = []
        for i, ev in enumerate(live_response['evidence'][:2], 1):
            evidence_text.append(
                f"{i}. {ev['diagnosis_label']}: runtime support returned by the backend"
            )

        add_bullet_box(
            slide, LEFT_COL_X, ROW4_Y, LEFT_COL_WIDTH + GUTTER + RIGHT_COL_WIDTH, ROW4_HEIGHT,
        "Runtime Support Returned by Demo KB",
            evidence_text,
            LIGHT_BLUE, BLUE, bullet_size=8
        )

    add_footer(slide, slide_num)


def add_full_input_appendix_slide(prs: Presentation, case_id: str, slide_num: int) -> None:
    result = load_case_result(case_id)
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    add_title(slide, f"Full Blinded Input: {case_id}", "Appendix case evidence")

    add_label(slide, 0.55, 1.25, 1.85, 0.48, "Case ID", case_id, LIGHT_TEAL, TEAL)
    add_label(slide, 2.55, 1.25, 1.6, 0.48, "Input mode", "Text + image", LIGHT_BLUE, BLUE)
    add_label(slide, 4.3, 1.25, 2.25, 0.48, "Source", "Held-out query", LIGHT_GREEN, GREEN)
    add_label(slide, 6.7, 1.25, 2.5, 0.48, "Use", "Q&A appendix", LIGHT_ORANGE, ORANGE)

    clinical_text = result["clinical_text"]
    content_y = 1.80  # Start slightly higher for more space

    # Adaptive layout based on text length - imported from v11b with improvements
    if len(clinical_text) > 1400:
        # Very long text: use two-column layout (v11b's superior approach)
        paragraphs = clinical_text.split("\n\n")
        left_body = "\n\n".join(paragraphs[:2])
        right_body = "\n\n".join(paragraphs[2:])

        add_text_panel(
            slide,
            0.55,
            content_y,
            4.24,  # Left column width
            3.20,
            "Full blinded clinical input",
            left_body,
            LIGHT_BLUE,
            BLUE,
            body_size=8.85,
        )
        add_text_panel(
            slide,
            4.96,  # Right column starts here
            content_y,
            4.24,  # Right column width
            3.20,
            "Full blinded clinical input (continued)",
            right_body,
            LIGHT_BLUE,
            BLUE,
            body_size=8.85,
        )
        image_box = None  # No image in two-column mode
    elif len(clinical_text) > 950:
        # Medium-long text: text + smaller image
        text_width = 6.22
        text_height = 3.20
        body_size = 8.75
        image_box = (6.98, content_y + 0.48, 1.95, 1.95)
    else:
        # Short text: text + larger image
        text_width = 5.55
        text_height = 3.20
        body_size = 9.8
        image_box = (6.38, content_y + 0.38, 2.55, 2.15)

    # Add single text panel for non-two-column layouts
    if len(clinical_text) <= 1400:
        add_text_panel(
            slide,
            0.55,
            content_y,
            text_width,
            text_height,
            "Full blinded clinical input",
            clinical_text,
            LIGHT_BLUE,
            BLUE,
            body_size=body_size,
        )

    # Add image if layout includes it
    image_path = EVAL_DIR / f"{case_id}_input.png"
    if image_box and image_path.exists():
        image_x, image_y, image_w, image_h = image_box
        add_picture_fit(slide, image_path, image_x, image_y, image_w, image_h)
        add_textbox(slide, image_x, image_y + image_h + 0.08,
                    image_w, 0.22, "Input image used in test", 7.5, SLATE, bold=True)

    add_textbox(
        slide,
        0.55,
        5.02,
        8.35,
        0.2,
        "Full blinded clinical input: diagnosis, title, abstract, captions, and confirmatory findings withheld.",
        7.4,
        SLATE,
    )
    add_footer(slide, slide_num)


def load_comparison_summary() -> dict:
    if not COMPARISON_SUMMARY.exists():
        raise FileNotFoundError(
            "Gemma 4 RAG/no-RAG comparison artifact missing. "
            "Run scripts/extract_gemma4_rag_norag_comparison.py first."
        )
    with COMPARISON_SUMMARY.open("r", encoding="utf-8") as handle:
        summary = json.load(handle)
    if summary.get("comparison_label") != "official Gemma 4 experiment-pipeline comparison":
        raise ValueError("Unexpected comparison source label")
    return summary


def load_trace_summary() -> dict:
    if not TRACE_SUMMARY.exists():
        raise FileNotFoundError(
            "Exact RAG trace appendix artifact missing. "
            "Run scripts/extract_exact_rag_trace_appendix.py first."
        )
    with TRACE_SUMMARY.open("r", encoding="utf-8") as handle:
        summary = json.load(handle)
    if summary.get("version") != "V12d":
        raise ValueError("Unexpected exact trace appendix version")
    return attach_live_recapture_audits(summary)


def short_rank(text: str, max_chars: int = 74) -> str:
    text = " ".join(text.split())
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "..."


def exact_prefix(text: str, max_chars: int) -> str:
    text = " ".join(str(text).split())
    if len(text) <= max_chars:
        return text
    return text[:max_chars].rstrip()


def score_text(value) -> str:
    return str(value)


def add_condition_box(slide, x, y, w, h, title, result, fill, accent, source_label):
    bullets = [
        f"Rank 1 type: {result['diagnosis_type']}",
        f"Rank 1: {short_rank(result['rank1'])}",
        f"Confidence: {result['confidence']}",
        f"Retrieved context chunks: {result['retrieved_context_count']}",
        source_label,
    ]
    add_bullet_box(slide, x, y, w, h, title, bullets, fill, accent, bullet_size=8.3)


def add_rag_norag_comparison_slide(
    prs: Presentation,
    comparison_summary: dict,
    case: dict,
    slide_num: int,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    case_id = case["case_id"]
    add_title(
        slide,
        f"Gemma 4 RAG/No-RAG Backup: {case_id}",
        "Appendix experiment-pipeline evidence",
    )

    add_label(slide, 0.55, 1.25, 1.95, 0.5, "Case ID", case_id, LIGHT_TEAL, TEAL)
    add_label(
        slide,
        2.65,
        1.25,
        1.85,
        0.5,
        "Silver reference",
        case["silver_reference_type"],
        LIGHT_BLUE,
        BLUE,
    )
    add_label(
        slide,
        4.65,
        1.25,
        2.1,
        0.5,
        "Use",
        "Q&A backup only",
        LIGHT_ORANGE,
        ORANGE,
    )
    add_label(
        slide,
        6.9,
        1.25,
        2.3,
        0.5,
        "Source",
        "Gemma 4 pipeline",
        LIGHT_GREEN,
        GREEN,
    )

    source_text = "official Gemma 4 experiment-pipeline comparison"
    add_textbox(
        slide,
        0.55,
        1.86,
        8.65,
        0.24,
        f"Source: {source_text}; same held-out case ID, Q1/Q3 multimodal query type.",
        8.5,
        SLATE,
    )

    add_condition_box(
        slide,
        0.55,
        2.22,
        4.2,
        1.55,
        "NO-RAG condition",
        case["no_rag"],
        LIGHT_RED,
        RED,
        "Experiment pipeline: no retrieved context",
    )
    add_condition_box(
        slide,
        5.0,
        2.22,
        4.2,
        1.55,
        "RAG condition",
        case["rag"],
        LIGHT_GREEN,
        GREEN,
        "Experiment pipeline: retrieved context used",
    )

    add_bullet_box(
        slide,
        0.55,
        4.00,
        8.65,
        1.05,
        "Presentation interpretation",
        [
            case["interpretation"],
            "Selected case-level backup only; aggregate thesis metrics remain the benchmark.",
            "Not clinical validation and not diagnosis from image alone.",
        ],
        LIGHT_BLUE,
        BLUE,
        bullet_size=8.6,
    )

    add_comparison_footer(slide, slide_num)


def add_context_trace_box(slide, x, y, w, h, context: dict) -> None:
    body = "\n".join(
        [
            f"doc_id: {context['doc_id']}",
            f"score: {score_text(context['score'])}",
            f"diagnosis_type: {context['diagnosis_type']} | label_source: {context['label_source']}",
            f"text_prefix: {exact_prefix(context['text_prefix_260'], 230)}",
        ]
    )
    add_text_panel(
        slide,
        x,
        y,
        w,
        h,
        f"Context rank {context['rank']}",
        body,
        LIGHT_BLUE,
        BLUE,
        body_size=5.8,
    )


def add_official_trace_slide(
    prs: Presentation,
    trace_summary: dict,
    case: dict,
    slide_num: int,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    case_id = case["case_id"]
    official = case["official_rag_trace"]
    run = trace_summary["official_rag_run"]

    add_title(
        slide,
        f"Official RAG Trace: {case_id}",
        "Appendix retriever/reranker evidence",
    )
    add_label(slide, 0.55, 1.25, 1.65, 0.48, "Case ID", case_id, LIGHT_TEAL, TEAL)
    add_label(
        slide,
        2.35,
        1.25,
        1.55,
        0.48,
        "Retriever",
        str(run["retriever_method"]),
        LIGHT_BLUE,
        BLUE,
    )
    add_label(
        slide,
        4.05,
        1.25,
        1.35,
        0.48,
        "Rerank",
        str(run["rerank"]),
        LIGHT_GREEN,
        GREEN,
    )
    add_label(
        slide,
        5.55,
        1.25,
        1.35,
        0.48,
        "Top K",
        str(run["retrieval_top_k"]),
        LIGHT_BLUE,
        BLUE,
    )
    add_label(
        slide,
        7.05,
        1.25,
        2.15,
        0.48,
        "Use",
        "Q&A trace only",
        LIGHT_ORANGE,
        ORANGE,
    )
    add_textbox(
        slide,
        0.55,
        1.83,
        8.65,
        0.22,
        f"qid: {case['qid']} | source: official Gemma 4 experiment-pipeline trace",
        7.8,
        SLATE,
    )
    add_textbox(
        slide,
        0.55,
        2.08,
        8.65,
        0.20,
        "Visible contexts are the rerank-enabled final context list used for generation; no separate pre-rerank list is claimed.",
        7.3,
        SLATE,
    )

    contexts = official["top_contexts_for_slide"][:3]
    for x, context in zip((0.55, 3.47, 6.38), contexts):
        add_context_trace_box(slide, x, 2.42, 2.82, 2.55, context)

    add_trace_footer(slide, slide_num)


def add_fresh_gpu_audit_slide(
    prs: Presentation,
    case: dict,
    slide_num: int,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    case_id = case["case_id"]
    audit = case["fresh_gpu_audit"]
    fields = audit["fields"]

    add_title(
        slide,
        f"Fresh Real-GPU Output Audit: {case_id}",
        "Appendix exact backend evidence",
    )
    add_label(slide, 0.55, 1.25, 1.8, 0.48, "Case ID", case_id, LIGHT_TEAL, TEAL)
    add_label(
        slide,
        2.5,
        1.25,
        2.05,
        0.48,
        "Provider",
        fields["provider_mode"],
        LIGHT_GREEN,
        GREEN,
    )
    add_label(
        slide,
        4.7,
        1.25,
        2.35,
        0.48,
        "Model",
        fields["model_name"],
        LIGHT_BLUE,
        BLUE,
    )
    add_label(
        slide,
        7.2,
        1.25,
        2.0,
        0.48,
        "Use",
        "Fresh audit only",
        LIGHT_ORANGE,
        ORANGE,
    )

    add_textbox(
        slide,
        0.55,
        1.85,
        8.65,
        0.23,
        f"request_id: {fields['request_id']} | source: {audit['source_path']}",
        7.0,
        SLATE,
    )

    result_bullets = [
        f"Rank 1 supportive consideration: {fields['rank1']}",
        f"Rank 1 Diagnosis Type: {fields['rank1_type']}",
        f"Rank 1 Confidence: {fields['rank1_confidence']}",
        f"safety_state: {fields['safety_state']} | elapsed_seconds: {fields['elapsed_seconds']} | query_image_tensor_count: {fields['query_image_tensor_count']}",
    ]
    add_bullet_box(
        slide,
        0.55,
        2.18,
        8.65,
        1.02,
        "Exact model-output fields from fresh backend response",
        result_bullets,
        LIGHT_GREEN,
        GREEN,
        bullet_size=7.2,
    )

    add_textbox(
        slide,
        0.55,
        3.34,
        8.65,
        0.2,
        "Exact returned evidence chunks from raw backend response",
        9.0,
        BLUE,
        True,
        PP_ALIGN.CENTER,
    )

    for x, evidence in zip((0.55, 3.47, 6.38), audit["evidence_for_slide"][:3]):
        body = "\n".join(
            [
                f"chunk_id: {evidence['chunk_id']} | score: {score_text(evidence['score'])}",
                f"title: {evidence['title']}",
                f"diagnosis_label: {evidence['diagnosis_label']} | confirmatory: {evidence['confirmatory']}",
                f"excerpt: {exact_prefix(evidence['excerpt'], 210)}",
            ]
        )
        add_text_panel(
            slide,
            x,
            3.62,
            2.82,
            1.48,
            f"Evidence rank {evidence['rank']}",
            body,
            LIGHT_TEAL,
            TEAL,
            body_size=5.55,
        )

    add_trace_footer(slide, slide_num)


def main() -> None:
    if not SOURCE_PPTX.exists():
        raise FileNotFoundError(f"Source PPTX missing: {SOURCE_PPTX}")

    if not (EVAL_DIR / "evaluation_summary.json").exists():
        raise FileNotFoundError(
            "Held-out evaluation results missing. "
            "Run scripts/run_heldout_case_evaluation.py first."
        )

    if not LOGO_PATH.exists():
        raise FileNotFoundError(f"Institutional logo missing: {LOGO_PATH}")
    comparison_summary = load_comparison_summary()
    trace_summary = load_trace_summary()

    prs = Presentation(str(SOURCE_PPTX))
    print(f"Loaded source presentation with {len(prs.slides)} slides")
    refine_core_wording(prs)
    print("Refined slides 12-14 for defense-facing wording")

    # Add overview slide
    add_case_overview_slide(prs, 15)
    print("Added slide 15: Overview")

    # Add case result slides
    add_case_result_slide(prs, "PMC7516301_01", 16, "success")
    print("Added slide 16: Case 1 (MCL concordant example)")

    add_case_result_slide(prs, "PMC7456484_01", 17, "limitation")
    print("Added slide 17: Case 2 (PKDL subtype limitation)")

    add_case_result_slide(prs, "PMC10026180_04", 18, "label_conflict")
    print("Added slide 18: Case 3 (label-conflict stress test)")

    add_closing_slide(prs)
    print("Added slide 19: Thank You / Questions")

    add_full_input_appendix_slide(prs, "PMC7516301_01", 20)
    print("Added slide 20: Full input appendix (PMC7516301_01)")

    add_full_input_appendix_slide(prs, "PMC7456484_01", 21)
    print("Added slide 21: Full input appendix (PMC7456484_01)")

    add_full_input_appendix_slide(prs, "PMC10026180_04", 22)
    print("Added slide 22: Full input appendix (PMC10026180_04)")

    for slide_num, case in enumerate(comparison_summary["cases"], start=23):
        add_rag_norag_comparison_slide(prs, comparison_summary, case, slide_num)
        print(f"Added slide {slide_num}: Gemma 4 RAG/no-RAG backup ({case['case_id']})")

    trace_slide_num = 26
    for case in trace_summary["cases"]:
        add_official_trace_slide(prs, trace_summary, case, trace_slide_num)
        print(f"Added slide {trace_slide_num}: official RAG trace ({case['case_id']})")
        trace_slide_num += 1
        add_fresh_gpu_audit_slide(prs, case, trace_slide_num)
        print(f"Added slide {trace_slide_num}: fresh real-GPU audit ({case['case_id']})")
        trace_slide_num += 1

    prs.save(str(OUTPUT_PPTX))
    print(f"\n✓ Built: {OUTPUT_PPTX}")
    print(f"  Total slides: {len(prs.slides)}")
    print(f"  Original: 14 slides")
    print(
        "  Added: 4 evaluation slides (15-18) + closing slide (19) "
        "+ 3 input appendix slides (20-22) + 3 comparison appendix slides (23-25) "
        "+ 6 exact trace/audit slides (26-31)"
    )


if __name__ == "__main__":
    main()
